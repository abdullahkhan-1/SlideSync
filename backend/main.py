from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import uvicorn
import fitz  # PyMuPDF
from pptx import Presentation
import io
import re

app = FastAPI(
    title="SlideSync API",
    docs_url="/docs",  # we'll disable this in production
    redoc_url=None
)

# Allow frontend to talk to backend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─── Helper: Extract from PDF ───────────────────────────────────────────────

def extract_from_pdf(file_bytes: bytes):
    slides = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict")["blocks"]
        
        title = ""
        content_lines = []
        max_font_size = 0

        for block in blocks:
            if block["type"] == 0:  # text block
                for line in block["lines"]:
                    for span in line["spans"]:
                        text = span["text"].strip()
                        font_size = span["size"]
                        
                        if not text:
                            continue
                        
                        # Largest font on page is likely the title
                        if font_size > max_font_size:
                            max_font_size = font_size
                            title = text
                        else:
                            content_lines.append(text)

        # Fallback: if no title found, use first line
        if not title and content_lines:
            title = content_lines.pop(0)

        content = " ".join(content_lines)

        slides.append({
            "slide_number": page_num,
            "title": title if title else f"Slide {page_num}",
            "content": content[:500]  # limit content length
        })

    doc.close()
    return slides

# ─── Helper: Extract from PPTX ──────────────────────────────────────────────

def extract_from_pptx(file_bytes: bytes):
    slides = []
    prs = Presentation(io.BytesIO(file_bytes))

    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        content_lines = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Slide title placeholder
                if shape.shape_type == 13:
                    continue
                    
                if not title:
                    # Check if this shape is a title placeholder
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0 or ph_idx == 1:
                            title = text
                            continue

                content_lines.append(text)

        # Fallback: use first content line as title
        if not title and content_lines:
            title = content_lines.pop(0)

        content = " ".join(content_lines)

        slides.append({
            "slide_number": slide_num,
            "title": title if title else f"Slide {slide_num}",
            "content": content[:500]
        })

    return slides

# ─── Main Endpoint ───────────────────────────────────────────────────────────

@app.post("/api/extract")
async def extract_slides(file: UploadFile = File(...)):
    
    # Validate file type
    filename = file.filename.lower()
    if not (filename.endswith(".pdf") or filename.endswith(".pptx")):
        raise HTTPException(
            status_code=400,
            detail="Only PDF and PPTX files are supported"
        )

    # Read file bytes
    file_bytes = await file.read()

    if len(file_bytes) == 0:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    # Extract based on type
    try:
        if filename.endswith(".pdf"):
            slides = extract_from_pdf(file_bytes)
        else:
            slides = extract_from_pptx(file_bytes)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

    return {
        "filename": file.filename,
        "total_slides": len(slides),
        "slides": slides
    }

# ─── Health Check ────────────────────────────────────────────────────────────

@app.get("/")
def health_check():
    return {"status": "SlideSync backend is running"}

# ─── Run ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)